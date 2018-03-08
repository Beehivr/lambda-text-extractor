from datetime import datetime, time
import io
import json
import logging
import os
import re
from subprocess import CalledProcessError
from tempfile import NamedTemporaryFile
from zipfile import is_zipfile, ZipFile

import boto3
import botocore

from utils import get_subprocess_output
from utils import get_ext
from uriutils import uri_read, uri_exists, uri_dump

LAMBDA_TASK_ROOT = os.environ.get('LAMBDA_TASK_ROOT', os.path.dirname(os.path.abspath(__file__)))
BIN_DIR = os.path.join(LAMBDA_TASK_ROOT, 'bin')
LIB_DIR = os.path.join(LAMBDA_TASK_ROOT, 'lib')

# TEXTRACTOR_OCR = os.environ['TEXTRACTOR_OCR']

with NamedTemporaryFile(mode='w', delete=False) as f:
    CATDOCRC_PATH = f.name
    f.write('charset_path = {}\n'.format(os.path.join(LIB_DIR, 'catdoc', 'charsets')))
    f.write('map_path = {}\n'.format(os.path.join(LIB_DIR, 'catdoc', 'charsets')))
#end with

logging.basicConfig(format='%(asctime)-15s [%(name)s-%(process)d] %(levelname)s: %(message)s', level=logging.INFO)

__name__ = 'simple_extractor'
logger = logging.getLogger(__name__)



def extract(event_handler, ext=None, disable_ocr=True):

    logger.info("\n\n\nStarting Document parse\n" + event_handler.key + "\n\n\n")

    # AWS Lambda auto-retries errors for 3x. This should make it disable retrying...kinda. See https://stackoverflow.com/questions/32064038/aws-lambda-function-triggering-multiple-times-for-a-single-event
    #aws_context_retry_uri = os.path.join(temp_uri_prefix, 'aws_lambda_request_ids', context.aws_request_id)
    #if uri_exists(aws_context_retry_uri):
    #    return
    #uri_dump(aws_context_retry_uri, '', mode='w')

    # logger.info('{} invoked with event {}.'.format(os.environ['AWS_LAMBDA_FUNCTION_NAME'], json.dumps(event)))
    if ext is None:
        ext = get_ext(event_handler.key)
    logger.info("file ext " + ext)

    extract_func = PARSE_FUNCS.get(ext)
    if extract_func is None:
        raise ValueError('<{}> has unsupported extension "{}".'.format(event_handler.key, ext))
        return

    fallback_to_ocr = False
    # textractor_results = {}
    if extract_func is False:
        fallback_to_ocr = True
        logger.info('Fallback to OCR for <{}>.'.format(event_handler.key))
        return ('ocr', None)

    else:
        with NamedTemporaryFile(mode='wb', suffix=ext, delete=False) as f:
            document_path = f.name
            f.write(event_handler.get_content().read())

            logger.info('Downloaded <{}> to <{}>.'.format(event_handler.key, document_path))
        #end with

        text = extract_func(document_path)

        if extract_func is pdf_to_text and len(text) < 512 and not disable_ocr:
            return ('ocr', None)

        else:
            if len(text) == 0: logger.warning('<{}> does not contain any content.'.format(event_handler.key))
            return ('ok', text)
        #end if

    #end if

    return ('ocr', None)

#end def


def _get_subprocess_output(*args, **kwargs):
    global logger
    kwargs['logger'] = logger
    return get_subprocess_output(*args, **kwargs)
#end def


#================================================================
#
#
#                        Text converters
#
#
#================================================================


def pdf_to_text(document_path):
    logger.info("\n>>> Starting pdf_to_text > " + document_path)
    with NamedTemporaryFile(suffix='.txt', delete=False) as f:
        text_path = f.name

    logger.info(">>> _get_subprocess_output ")
    _get_subprocess_output([os.path.join(BIN_DIR, 'pdftotext'), '-layout', '-nopgbrk', '-eol', 'unix', document_path, text_path], shell=False, env=dict(LD_LIBRARY_PATH=os.path.join(LIB_DIR, 'pdftotext')))

    logger.info(">>> io_open ")
    with io.open(text_path, mode='r', encoding='utf-8', errors='ignore') as f:
        text = f.read().strip()

    logger.info(">>> os_remove ")
    # os.remove(text_path)

    logger.info(">>> done ")

    return text
#end def


def doc_to_text(document_path, event_handler):
    global logger

    cmdline = [os.path.join(BIN_DIR, 'antiword'), '-t', '-w', '0', '-m', 'UTF-8', document_path]
    try:
        text = _get_subprocess_output(cmdline, display_output_on_exception=False, shell=False, env=dict(ANTIWORDHOME=os.path.join(LIB_DIR, 'antiword')))
        text = text.decode('utf-8', errors='ignore').strip()
    except CalledProcessError as e:
        if b'Rich Text Format' in e.output:
            logger.info('Antiword failed on possible Rich Text file <{}>.'.format(event_handler.key))
            return rtf_to_text(document_path, event_handler)

        elif b'"docx" file' in e.output or is_zipfile(document_path):
            logger.info('Antiword failed on possible docx file <{}>.'.format(event_handler.key))
            return docx_to_text(document_path, event_handler)

        else:
            logger.exception('Antiword exception with output "{}".'.format(e.output.decode('ascii', errors='ignore')))
            text = None
        #end if
    #end try

    if text is None:  # Fallback to catdoc
        cmdline = [os.path.join(BIN_DIR, 'catdoc'), '-a', document_path]
        text = _get_subprocess_output(cmdline, shell=False, env=dict(CATDOCRC_PATH=CATDOCRC_PATH))
        text = text.decode('utf-8', errors='ignore').strip()
    #end if

    return text
#end def


def docx_to_text(document_path, event_handler):
    global logger

    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    try:
        doc = Document(document_path)
        doc_body = doc.element.body
        blocks = []
        for child in doc_body.iterchildren():
            if isinstance(child, CT_P): blocks.append(Paragraph(child, doc_body).text)
            elif isinstance(child, CT_Tbl): blocks.append('\n'.join(' | '.join(cell.text for cell in row.cells) for row in Table(child, doc_body).rows))
        #end for

        text = '\n\n'.join(blocks).strip()

        return text

    except Exception:
        logger.exception('Exception while parsing <{}>.'.format(event_handler.key))
    #end try

    # Extract it from the XML
    with ZipFile(document_path) as document_zipfile:
        xml_content = document_zipfile.read('word/document.xml')

    try: from xml.etree.cElementTree import XML
    except ImportError: from xml.etree.ElementTree import XML

    tree = XML(xml_content)

    DOCX_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    DOCX_PARA = DOCX_NAMESPACE + 'p'
    DOCX_TEXT = DOCX_NAMESPACE + 't'

    paragraphs = []
    for paragraph in tree.getiterator(DOCX_PARA):
        texts = [node.text for node in paragraph.getiterator(DOCX_TEXT) if node.text]
        if texts:
            paragraphs.append(''.join(texts))
    #end for

    text = '\n\n'.join(paragraphs)

    return text
#end def


def rtf_to_text(document_path):
    cmdline = [os.path.join(BIN_DIR, 'unrtf'), '-P', os.path.join(LIB_DIR, 'unrtf'), '--text', document_path]
    text = _get_subprocess_output(cmdline, shell=False)
    text = text.decode('utf-8', errors='ignore')

    new_lines = []
    in_header = True
    for line in text.split('\n'):
        if in_header and line.startswith('###'): continue
        else:
            new_lines.append(line)
            in_header = False
        #end if
    #end for
    text = '\n'.join(new_lines).strip()
    text = re.sub(r'[\x0e-\x1f]', '', text)

    return text
#end def


def xls_to_text(document_path):
    import xlrd

    book = xlrd.open_workbook(document_path)
    lines = []
    for sheetno, sheet in enumerate(book.sheets()):
        lines.append(sheet.name)
        lines.append('-------------------------------------------')
        for row in sheet.get_rows():
            row_values = []
            for cell in row:
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try: d = datetime(*xlrd.xldate_as_tuple(cell.value, book.datemode))
                    except ValueError: d = datetime(1970, 1, 1)
                    row_values.append(d.date() if d.time() == time(0, 0) else d)
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN: row_values.append(bool(cell.value))
                else: row_values.append(cell.value)
            #end for
            lines.append(' | '.join(map(lambda s: str(s).strip(), row_values)))
        #end for
        lines.append('')  # empty line
    #end for

    return '\n'.join(lines).strip()
#end def


# def ppt_to_text(document_path, event, context):
#     cmdline = [os.path.join(BIN_DIR, 'catppt'), '-dutf-8', document_path]
#     text = _get_subprocess_output(cmdline, shell=False, env=dict(CATDOCRC_PATH=CATDOCRC_PATH))

#     return text.decode('utf-8', errors='ignore').strip()
# #end def


def pptx_to_text(document_path):
    import pptx

    prs = pptx.Presentation(document_path)

    paragraphs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            paragraphs += [' '.join([run.text for run in paragraph.runs]) for paragraph in shape.text_frame.paragraphs]
        #end for
    #end for

    return '\n\n'.join(paragraphs).strip()
#end def


def html_to_text(document_path):
    import lxml.html

    document = lxml.html.parse(document_path)
    for tag in ['script', 'style']:
        for elem in document.xpath('//' + tag):
            elem.drop_tree()

    return '\n'.join(document.xpath("//*//text()")).strip()
#end def


def text_to_text(document_path):
    with io.open(document_path, mode='r', encoding='utf-8', errors='ignore') as f:
        text = f.read().strip()

    return text
#end def


def csv_to_text(document_path):
    import csv
    with io.open(document_path, mode='r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        text = '\n'.join('\t'.join(row) for row in reader)
    #end with

    return text
#end def


def odf_to_text(document_path):
    from odf.opendocument import load as odf_load
    from odf import text as odf_text
    from odf import teletype as odf_teletype

    doc = odf_load(document_path)
    paragraphs = []
    for p in doc.getElementsByType(odf_text.P):
        paragraphs.append(odf_teletype.extractText(p))

    return '\n'.join(paragraphs).strip()
#end def


PARSE_FUNCS = {
    '.csv': csv_to_text,
    '.doc': doc_to_text,
    '.docx': docx_to_text,
    '.dot': doc_to_text,
    '.htm': html_to_text,
    '.html': html_to_text,
    '.odm': odf_to_text,
    '.odp': odf_to_text,
    '.ods': odf_to_text,
    '.odt': odf_to_text,
    '.oth': odf_to_text,  # not tested
    '.otm': odf_to_text,
    '.otp': odf_to_text,
    '.ots': odf_to_text,
    '.ott': odf_to_text,
    '.pdf': pdf_to_text,
    # '.ppt': ppt_to_text,  # catppt not working (2017/08/10)
    '.pptx': pptx_to_text,
    '.rtf': rtf_to_text,
    '.text': text_to_text,
    '.txt': text_to_text,
    '.xls': xls_to_text,
    '.xlsx': xls_to_text,
    # image ones
    '.png': False,  # False means default to OCR
    '.tiff': False,
    '.tif': False,
    '.jpg': False,
    '.jpeg': False,
}
